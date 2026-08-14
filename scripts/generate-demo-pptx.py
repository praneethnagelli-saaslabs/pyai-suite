#!/usr/bin/env python3
"""Generate docs/PyAI-Suite-Demo.pptx for hackathon demos."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "docs" / "PyAI-Suite-Demo.pptx"

# Charcoal + teal — avoid purple / cream clichés
BG = RGBColor(0x12, 0x16, 0x1C)
FG = RGBColor(0xE8, 0xEC, 0xF1)
MUTED = RGBColor(0x9A, 0xA3, 0xB0)
ACCENT = RGBColor(0x2F, 0xA3, 0x8A)


def set_slide_bg(slide, prs):
    fill = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        prs.slide_width,
        prs.slide_height,
    )
    fill.fill.solid()
    fill.fill.fore_color.rgb = BG
    fill.line.fill.background()
    spTree = slide.shapes._spTree
    sp = fill._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_text(slide, left, top, width, height, text, *, size=28, bold=False, color=FG, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Helvetica Neue"
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, lines, *, size=20, color=FG):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Helvetica Neue"
        p.space_after = Pt(10)
    return box


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, prs)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Title
    s = new_slide(prs)
    add_text(s, 0.8, 1.8, 11, 0.5, "SAAS LABS", size=16, color=ACCENT, bold=True)
    add_text(s, 0.8, 2.3, 11, 1.2, "PyAI Suite", size=54, bold=True)
    add_text(
        s,
        0.8,
        3.6,
        11,
        1,
        "Four AI products on one voice platform — PyAI first, with clear fallbacks.",
        size=22,
        color=MUTED,
    )
    add_text(s, 0.8, 5.8, 11, 0.5, "CallIQ  ·  Scrib  ·  Brief  ·  Simulator", size=18, color=ACCENT)

    # 2 Built vs pulled in
    s = new_slide(prs)
    add_text(s, 0.8, 0.5, 11, 0.6, "Before demos — own this line", size=16, color=ACCENT, bold=True)
    add_text(s, 0.8, 1.0, 11, 0.8, "Built vs pulled in", size=40, bold=True)
    add_text(s, 0.8, 2.0, 5.5, 0.4, "WE BUILT", size=14, color=ACCENT, bold=True)
    add_bullets(
        s,
        0.8,
        2.5,
        5.5,
        4,
        [
            "CallIQ · Scrib · Brief · Simulator",
            "Provider registry & visible fallbacks",
            "API, web UI, extension, desktop tray",
            "Recap, cleanup, meeting memory, scoring",
        ],
        size=18,
    )
    add_text(s, 7.0, 2.0, 5.5, 0.4, "WE PULLED IN", size=14, color=ACCENT, bold=True)
    add_bullets(
        s,
        7.0,
        2.5,
        5.5,
        4,
        [
            "PyAI Hear / Speak",
            "OpenAI / Gemini adapters",
            "Attendee (open-source Meet bot)",
            "Postgres · Redis · MinIO · Docker",
        ],
        size=18,
    )

    # 3 One-liners
    s = new_slide(prs)
    add_text(s, 0.8, 0.5, 11, 0.6, "Products", size=16, color=ACCENT, bold=True)
    add_text(s, 0.8, 1.0, 11, 0.8, "One line each", size=40, bold=True)
    rows = [
        ("CallIQ", "Listens to a sales call → evidence-backed deal notes"),
        ("Scrib", "You speak → cleaned text ready to paste"),
        ("Brief", "Listens to a meeting → notes + searchable memory"),
        ("Simulator", "Practice / stress-test a live voice agent"),
    ]
    y = 2.2
    for title, body in rows:
        add_text(s, 0.8, y, 2.5, 0.45, title, size=22, bold=True, color=ACCENT)
        add_text(s, 3.5, y, 9, 0.45, body, size=20, color=FG)
        y += 0.85

    # 4 Stack
    s = new_slide(prs)
    add_text(s, 0.8, 0.5, 11, 0.6, "Tech stack", size=16, color=ACCENT, bold=True)
    add_text(s, 0.8, 1.0, 11, 0.8, "What powers the demos", size=40, bold=True)
    add_bullets(
        s,
        0.8,
        2.2,
        11,
        4.5,
        [
            "PyAI — Hear (STT) · Speak (TTS) · primary path",
            "Fallbacks — OpenAI · Gemini · Mock when no keys",
            "Attendee — CallIQ Bot joins Google Meet / Zoom",
            "Infra — Postgres · Redis · MinIO · Docker Compose",
            "Apps — React web · Fastify API · Chrome extension · Tauri tray",
        ],
        size=22,
    )

    # 5 Demo order
    s = new_slide(prs)
    add_text(s, 0.8, 0.5, 11, 0.6, "Live demo", size=16, color=ACCENT, bold=True)
    add_text(s, 0.8, 1.0, 11, 0.8, "Click path (~12–15 min)", size=40, bold=True)
    add_bullets(
        s,
        0.8,
        2.2,
        11,
        4.5,
        [
            "0  Built vs pulled in (20 sec)",
            "1  CallIQ → Run product demo (4–5 min)",
            "2  Scrib → Try demo (2 min)",
            "3  Brief → Try sample demo (2–3 min)",
            "4  Simulator → Start call (2 min)",
            "5  Providers → health + fallbacks (30 sec)",
            "Open http://localhost:3000",
        ],
        size=22,
    )

    # 6 CallIQ
    s = new_slide(prs)
    add_text(s, 0.8, 0.5, 11, 0.6, "FEATURED", size=16, color=ACCENT, bold=True)
    add_text(s, 0.8, 1.0, 11, 0.8, "CallIQ", size=44, bold=True)
    add_text(
        s,
        0.8,
        2.0,
        11,
        1,
        "After a sales call, someone still rewrites the CRM by hand.\nCallIQ listens → Hear transcripts → Recap writes deal notes.",
        size=22,
        color=MUTED,
    )
    add_bullets(
        s,
        0.8,
        3.5,
        11,
        3,
        [
            "Click: Run product demo",
            "Show: transcript · conversation shape · summary / risks / next steps",
            "Optional: Join Meet as bot (Attendee) — admit CallIQ Bot once",
        ],
        size=20,
    )

    # 7 Others
    s = new_slide(prs)
    add_text(s, 0.8, 0.5, 11, 0.6, "Also on the suite", size=16, color=ACCENT, bold=True)
    add_text(s, 0.8, 1.0, 11, 0.8, "Scrib · Brief · Simulator", size=40, bold=True)
    add_bullets(
        s,
        0.8,
        2.2,
        11,
        4.5,
        [
            "Scrib — Try demo · Speak → Hear → cleanup → paste-ready text",
            "Brief — Sample demo · or tab + mic capture → End meeting → notes + memory",
            "Simulator — Start a call · score the agent before customers hear it",
        ],
        size=22,
    )

    # 8 Failures
    s = new_slide(prs)
    add_text(s, 0.8, 0.5, 11, 0.6, "Stagecraft", size=16, color=ACCENT, bold=True)
    add_text(s, 0.8, 1.0, 11, 0.8, "If something fails", size=40, bold=True)
    add_bullets(
        s,
        0.8,
        2.2,
        11,
        4,
        [
            "Stay on Mock / sample demo paths",
            "Providers → healthy key ≠ Hear always wins (cooldown / fallback banners)",
            "Fallback banner = expected, not broken",
            "Skip live Meet — finish on CallIQ product demo",
        ],
        size=22,
    )

    # 9 Close
    s = new_slide(prs)
    add_text(s, 0.8, 2.2, 11, 1, "One voice platform. PyAI first.", size=36, bold=True, align=PP_ALIGN.CENTER)
    add_text(
        s,
        0.8,
        3.3,
        11,
        0.8,
        "CallIQ · Scrib · Brief · Simulator",
        size=22,
        color=ACCENT,
        align=PP_ALIGN.CENTER,
    )
    add_text(s, 0.8, 4.5, 11, 0.5, "Questions?", size=28, color=MUTED, align=PP_ALIGN.CENTER)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
